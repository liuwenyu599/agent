# -*- coding: utf-8 -*-
"""智能 PPT 引擎（移植自旧系统 ppt_service，AI 调用改经 WritingAssistant）。

核心设计：内容与版式分离
- 文档 outline 只保存内容：页面类型、标题、要点、结构化数据块、图片名；
- 模板（PPTTemplate）保存视觉风格（colors/font）+ 版式体系（layouts）。
  同一份内容换模板 = 换渲染样式，内容不动。

页型体系（10 种）：
  cover 封面 / toc 目录 / section 章节页 / content 正文（可选配图）
  data 数据卡片页 / chart 图表页（原生可编辑图表）
  case 案例页 / timeline 时间轴 / process 流程图 / summary 总结页 / closing 结束页

导出：python-pptx 手工排版，文本/形状/图表均为可编辑元素。
"""
import io
import json
import os
import re
import uuid
from typing import Dict, List, Optional

from app.core.logging import get_logger

logger = get_logger(__name__)

# ================= 内置官方模板（含完整版式体系） =================
BUILTIN_TEMPLATES = [
    {"id": "gov_report_red", "name": "政务工作汇报", "category": "工作汇报", "is_official": True,
     "description": "庄重政务红，适用：工作汇报、阶段总结",
     "colors": {"primary": "C00000", "accent": "E8B54D", "light": "FDF2F2", "dark": "7A0000"},
     "layouts": {"cover": "band_bottom", "toc": "numbered_list", "section": "left_block",
                  "content": "bar_title", "data": "cards_row", "chart": "chart_area",
                  "case": "header_text", "timeline": "horizontal", "process": "chevron",
                  "summary": "cards", "closing": "center"}},
    {"id": "policy_blue", "name": "政策解读", "category": "政策解读", "is_official": True,
     "description": "沉稳政策蓝，适用：政策解读、方案宣讲",
     "colors": {"primary": "1F4E9C", "accent": "5B9BD5", "light": "EEF4FC", "dark": "12325F"},
     "layouts": {"cover": "left_block", "toc": "numbered_list", "section": "center",
                  "content": "top_band", "data": "cards_row", "chart": "chart_area",
                  "case": "header_text", "timeline": "horizontal", "process": "chevron",
                  "summary": "cards", "closing": "brand_band"}},
    {"id": "training_blue", "name": "培训课件", "category": "培训课件", "is_official": True,
     "description": "明快培训蓝，适用：业务培训、知识讲座",
     "colors": {"primary": "2E75B6", "accent": "9DC3E6", "light": "F0F7FD", "dark": "1B4F82"},
     "layouts": {"cover": "band_bottom", "toc": "cards", "section": "left_block",
                  "content": "bar_title", "data": "cards_row", "chart": "chart_area",
                  "case": "header_text", "timeline": "horizontal", "process": "chevron",
                  "summary": "cards", "closing": "center"}},
    {"id": "exp_share_red", "name": "经验交流", "category": "经验交流", "is_official": True,
     "description": "热烈交流红，适用：经验交流、工作分享",
     "colors": {"primary": "B01E23", "accent": "F2C14E", "light": "FBF1F1", "dark": "6E1114"},
     "layouts": {"cover": "left_block", "toc": "cards", "section": "center",
                  "content": "bar_title", "data": "cards_row", "chart": "chart_area",
                  "case": "header_text", "timeline": "horizontal", "process": "chevron",
                  "summary": "cards", "closing": "center"}},
    {"id": "summary_brown", "name": "总结汇报", "category": "总结汇报", "is_official": True,
     "description": "厚重棕红，适用：年度总结、专项总结",
     "colors": {"primary": "8E1B1B", "accent": "D9A441", "light": "FAF3EE", "dark": "5C1010"},
     "layouts": {"cover": "band_bottom", "toc": "numbered_list", "section": "left_block",
                  "content": "top_band", "data": "cards_row", "chart": "chart_area",
                  "case": "header_text", "timeline": "horizontal", "process": "chevron",
                  "summary": "cards", "closing": "brand_band"}},
    {"id": "publicity_cyan", "name": "法治宣传", "category": "宣传展示", "is_official": True,
     "description": "清新宣传青，适用：法治宣传、普法教育",
     "colors": {"primary": "0F6E8C", "accent": "3FA7B8", "light": "ECF7F9", "dark": "0A4A5E"},
     "layouts": {"cover": "left_block", "toc": "cards", "section": "center",
                  "content": "bar_title", "data": "cards_row", "chart": "chart_area",
                  "case": "header_text", "timeline": "horizontal", "process": "chevron",
                  "summary": "cards", "closing": "center"}},
]

DEFAULT_LAYOUTS = BUILTIN_TEMPLATES[0]["layouts"]
DEFAULT_COLORS = BUILTIN_TEMPLATES[0]["colors"]
PAGE_TYPES = ["cover", "toc", "section", "content", "data", "chart",
              "case", "timeline", "process", "summary", "closing"]


def template_to_dict(t) -> Dict:
    """ORM 实体/字典 → 渲染用统一结构。"""
    if isinstance(t, dict):
        return {
            "id": t.get("id"), "name": t.get("name", "模板"), "category": t.get("category", "其他"),
            "colors": {**DEFAULT_COLORS, **(t.get("colors") or {})},
            "font": t.get("font") or "微软雅黑",
            "layouts": {**DEFAULT_LAYOUTS, **(t.get("layouts") or {})},
            "layout_library": t.get("layout_library") or [],
        }
    return {
        "id": t.id, "name": t.name, "category": t.category,
        "colors": {**DEFAULT_COLORS, **(t.colors or {})},
        "font": t.font or "微软雅黑",
        "layouts": {**DEFAULT_LAYOUTS, **(t.layouts or {})},
        "layout_library": t.layout_library or [],
    }


# ================= AI 大纲生成 =================

_OUTLINE_PROMPT = """你是司法行政机关的资深 PPT 策划专家。请根据材料设计演示文稿大纲。

【材料类型】{source_type}
【主题/要求】{topic}
【材料内容】
{content}

【用户素材库图片】（可按内容选用，选了就一定会在该页插入这张图片）
{image_list}

【页面类型】（按内容需要选用，不要全是 content）
- cover 封面（必有且仅第 1 页）：title、subtitle
- toc 目录页（超过 8 页时建议有）：points 为章节名（3~5 条）
- section 章节页：title
- content 正文页：title、points（3~5 条，每条 15~40 字的完整句子）
- data 数据页：title、blocks={{"cards":[{{"label":"指标名","value":"数值"}},...]}}（材料里有具体数据时用，3~4 张卡）
- chart 图表页：title、blocks={{"chart":{{"chart_type":"bar或pie或line","categories":["类别"...],"series":[{{"name":"系列名","values":[数值...]}}]}}}}（材料有对比/趋势数据时用）
- case 案例页：title、points（案例背景→做法→结果，3~4 条）
- timeline 时间轴页：title、blocks={{"timeline":[{{"time":"时间点","text":"事项"}},...]}}（有时间线内容时用）
- process 流程页：title、blocks={{"process":["步骤1","步骤2",...]}}（讲工作流程时用，3~5 步）
- summary 总结页：title、points（2~4 条结论/展望）
- closing 结束页（必有且仅最后一页）：title（如"谢谢聆听，请批评指正"）

【输出要求】
1. 只输出 JSON，不要用 markdown 代码块包裹，不要输出其他文字：
{{"title":"PPT主标题（≤20字）","subtitle":"副标题（≤30字）","slides":[{{"type":"cover","title":"...","subtitle":"..."}},...]}}
2. 每页 slide 的字段：type、title、points、blocks、image_name、note（演讲备注，一句话）；
3. 总页数 {slide_count} 页左右；
4. 内容要实：要点是有信息量的完整句子，材料中的具体数据、做法、案例必须用起来，禁止"加强领导""高度重视"式空话；
5. 图片规则：仅当某页内容与素材库图片明显相关时设 image_name（必须严格等于素材库中的图片名称），否则为 null；
6. 不要虚构材料中没有的事实和数据。"""


def generate_outline(assistant, source_type: str, topic: str, content: str,
                     images: List[Dict], slide_count: int = 10,
                     audience: str = "", scene: str = "") -> Dict:
    image_list = "\n".join(f"- {im['name']}：{im.get('caption') or '（无说明）'}" for im in images) \
        or "（素材库为空）"
    extra = ""
    if audience:
        extra += f"\n【汇报对象】{audience}（据此把握详略和措辞层级）"
    if scene:
        extra += f"\n【场景用途】{scene}"
    prompt = _OUTLINE_PROMPT.format(
        source_type=source_type, topic=(topic or "（未指定，依据材料自拟）") + extra,
        content=(content or "")[:6000], image_list=image_list,
        slide_count=slide_count,
    )
    try:
        raw = assistant.complete(
            [{"role": "system", "content": "你是 PPT 大纲策划模块，只输出 JSON。"},
             {"role": "user", "content": prompt}],
            temperature=0.5, max_tokens=4096,
        )
    except Exception as e:
        raise RuntimeError(f"调用模型失败: {e}")
    if raw.startswith("【系统错误】") or raw.startswith("调用模型失败"):
        raise RuntimeError(raw)
    m = re.search(r"\{.*\}", raw, re.S)
    if not m:
        raise RuntimeError("AI 未返回有效大纲，请重试")
    return normalize_outline(json.loads(m.group(0)))


def normalize_outline(outline: Dict) -> Dict:
    """规范化 + 兼容旧页型（bullets→content，image_text→content+配图版式）。"""
    type_alias = {"bullets": "content", "image_text": "content"}
    norm = []
    for s in (outline.get("slides") or []):
        t = s.get("type") or "content"
        legacy_image = (t == "image_text")
        t = type_alias.get(t, t)
        if t not in PAGE_TYPES:
            t = "content"
        layout = s.get("layout") or ("image_right" if legacy_image or s.get("image_name") else None)
        norm.append({
            "id": s.get("id") or uuid.uuid4().hex[:8],
            "type": t,
            "layout": layout,
            "title": (s.get("title") or "")[:60],
            "subtitle": (s.get("subtitle") or "")[:80],
            "points": [str(p)[:150] for p in (s.get("points") or [])][:6],
            "blocks": s.get("blocks") or {},
            "image_name": s.get("image_name") or None,
            "image_hint": s.get("image_hint") or "",
            "note": s.get("note") or "",
        })
    if not norm or norm[0]["type"] != "cover":
        norm.insert(0, {"id": uuid.uuid4().hex[:8], "type": "cover",
                        "title": outline.get("title", "演示文稿"),
                        "subtitle": outline.get("subtitle", ""), "layout": None,
                        "points": [], "blocks": {}, "image_name": None, "image_hint": "", "note": ""})
    if norm[-1]["type"] != "closing":
        norm.append({"id": uuid.uuid4().hex[:8], "type": "closing", "title": "谢谢聆听，请批评指正",
                     "subtitle": "", "layout": None, "points": [], "blocks": {},
                     "image_name": None, "image_hint": "", "note": ""})
    outline["slides"] = norm
    outline.setdefault("title", norm[0]["title"])
    outline.setdefault("subtitle", norm[0].get("subtitle", ""))
    return outline


def resolve_slide_images(outline: Dict, images: List[Dict]) -> Dict:
    """智能配图：AI 指定名 → 严格匹配 → 模糊匹配；找不到则不配图（版式自动退化）。"""
    by_name = {im["name"]: im for im in images}

    def _fuzzy(hint: str) -> Optional[Dict]:
        if not hint:
            return None
        hint_chars = set(re.sub(r"\s", "", hint.lower()))
        best, best_score = None, 0
        for im in images:
            score = len(hint_chars & set(im["name"].lower()))
            if score > best_score:
                best, best_score = im, score
        return best if best_score >= 2 else None

    for s in outline.get("slides", []):
        img = None
        if s.get("image_name"):
            img = by_name.get(s["image_name"]) or _fuzzy(s["image_name"])
        if not img and s.get("layout") == "image_right":
            img = _fuzzy(s.get("image_hint") or s.get("title") or "")
        s["_image"] = img
        if img and s.get("layout") is None and s.get("type") == "content":
            s["layout"] = "image_right"
    return outline


# ================= PPTX 构建 =================

def build_pptx(outline: Dict, template: Dict, base_pptx: bytes = None) -> bytes:
    """按模板渲染。template 为 template_to_dict 的结果。所有元素可编辑。

    base_pptx：用户上传的原始 pptx 字节。提供时以其为底版渲染——
    保留原模板的母版背景、装饰图形、主题配色，仅清空页面后逐页绘制内容。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    colors, layouts = template["colors"], template["layouts"]
    FONT = template.get("font") or "微软雅黑"
    C = lambda h: RGBColor.from_string(h)
    PRIMARY, ACCENT = C(colors["primary"]), C(colors["accent"])
    LIGHT, DARK = C(colors["light"]), C(colors["dark"])
    WHITE, GRAY, LINE = C("FFFFFF"), C("595959"), C("E7E6E6")

    if base_pptx:
        # 以用户上传的模板为底版：继承母版/主题/页面尺寸，清空原页面（连关系一起删，避免示例内容残留）
        prs = Presentation(io.BytesIO(base_pptx))
        sldIdLst = prs.slides._sldIdLst
        for sld in list(sldIdLst):
            rid = sld.get(qn('r:id'))
            sldIdLst.remove(sld)
            try:
                prs.part.drop_rel(rid)
            except Exception:
                pass
        # 选占位符最少的版式当空白页，避免母版占位符干扰自绘内容
        try:
            blank = min(prs.slide_layouts, key=lambda l: len(l.placeholders))
        except Exception:
            blank = prs.slide_layouts[0]
    else:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    # 底版母版为深色背景时，标题/正文文字自动改浅色，保证可读
    if base_pptx:
        try:
            mbg = prs.slide_masters[0].background.fill
            bgc = str(mbg.fore_color.rgb)
            br, bg_, bb = (int(bgc[i:i + 2], 16) for i in (0, 2, 4))
            if 0.299 * br + 0.587 * bg_ + 0.114 * bb < 110:
                GRAY = C("E6E6E6")
                DARK = C("FFFFFF")
        except Exception:
            pass

    def lx(f): return Emu(int(SW * f))
    def ly(f): return Emu(int(SH * f))

    def _font(run, size, bold=False, color=None):
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {})
            rPr.append(ea)
        ea.set('typeface', FONT)

    def _alpha(shape, pct):
        """给形状填充加透明度。pct=剩余不透明度（0~100）。"""
        spPr = shape._element.spPr
        fill = spPr.find(qn('a:solidFill'))
        if fill is None:
            return
        clr = fill.find(qn('a:srgbClr'))
        if clr is None:
            return
        a = clr.makeelement(qn('a:alpha'), {'val': str(int(pct * 1000))})
        clr.append(a)

    def rect(slide, x, y, w, h, color, round_=False, alpha=None):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        if round_:
            try:
                shp.adjustments[0] = 0.08
            except Exception:
                pass
        if alpha is not None:
            _alpha(shp, alpha)
        return shp

    def circle(slide, cx, cy, r, color, alpha):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - r)), Emu(int(cy - r)),
                                     Emu(int(2 * r)), Emu(int(2 * r)))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        _alpha(shp, alpha)
        return shp

    def text(slide, x, y, w, h, s, size, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = s
        _font(r, size, bold, color)
        return tb

    def picture_fit(slide, path, x, y, w, h):
        from PIL import Image
        with Image.open(path) as im:
            iw, ih = im.size
        box_ratio, img_ratio = w / h, iw / ih
        pic = slide.shapes.add_picture(path, x, y, w, h)
        if img_ratio > box_ratio:
            crop = (1 - box_ratio / img_ratio) / 2
            pic.crop_left = pic.crop_right = crop
        else:
            crop = (1 - img_ratio / box_ratio) / 2
            pic.crop_top = pic.crop_bottom = crop
        return pic

    def deco_circles(slide, light_on_dark=True):
        """封面/章节/结束页的装饰圆。"""
        base = ACCENT if light_on_dark else PRIMARY
        circle(slide, SW * 0.92, SH * 0.12, SH * 0.30, base, 18)
        circle(slide, SW * 0.98, SH * 0.55, SH * 0.18, base, 12)
        circle(slide, SW * 0.06, SH * 0.95, SH * 0.24, base, 14)

    def footer(slide, idx, total):
        """内容页页脚：细分隔线 + 页码。"""
        rect(slide, lx(0.045), ly(0.955), lx(0.91), Emu(9525), LINE)
        text(slide, lx(0.88), ly(0.962), lx(0.08), ly(0.03), f"{idx} / {total}", 9, False, GRAY,
             align=PP_ALIGN.RIGHT)

    def page_header(slide, title):
        """内容页页眉：bar_title=左侧竖条 / top_band=顶部色带。"""
        v = layouts.get("content", "bar_title")
        if v == "top_band":
            rect(slide, 0, 0, SW, ly(0.135), PRIMARY)
            rect(slide, 0, ly(0.135), SW, ly(0.006), ACCENT)
            text(slide, lx(0.05), ly(0.032), lx(0.9), ly(0.08), title, 24, True, WHITE)
        else:
            rect(slide, 0, 0, SW, ly(0.014), PRIMARY)
            rect(slide, lx(0.045), ly(0.07), lx(0.014), ly(0.082), PRIMARY, round_=True)
            text(slide, lx(0.072), ly(0.062), lx(0.85), ly(0.10), title, 26, True, DARK)
            rect(slide, lx(0.045), ly(0.175), lx(0.91), Emu(12700), LINE)

    def notes(slide, note):
        if note:
            slide.notes_slide.notes_text_frame.text = note

    def numbered_points(slide, pts, x=0.06, y=0.26, w=0.86, size=17):
        """编号要点：彩色序号圆角块 + 文字。"""
        for i, p in enumerate(pts):
            chip = rect(slide, lx(x), ly(y), lx(0.034), lx(0.034), PRIMARY, round_=True)
            tf = chip.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            r = para.add_run()
            r.text = str(i + 1)
            _font(r, 13, True, WHITE)
            text(slide, lx(x + 0.048), ly(y - 0.012), lx(w - 0.05), ly(0.12), p, size, False, GRAY)
            y += 0.128 if len(pts) <= 4 else 0.108
        return y

    # ---------- 页型渲染 ----------
    def r_cover(s):
        slide = prs.slides.add_slide(blank)
        v = layouts.get("cover", "band_bottom")
        if v == "left_block":
            rect(slide, 0, 0, lx(0.40), SH, PRIMARY)
            circle(slide, SW * 0.34, SH * 0.18, SH * 0.22, ACCENT, 18)
            circle(slide, SW * 0.08, SH * 0.88, SH * 0.26, ACCENT, 12)
            rect(slide, lx(0.40), 0, lx(0.008), SH, ACCENT)
            rect(slide, lx(0.50), ly(0.36), lx(0.05), ly(0.010), ACCENT)
            text(slide, lx(0.50), ly(0.40), lx(0.44), ly(0.24), s["title"], 36, True, DARK)
            if s.get("subtitle"):
                text(slide, lx(0.50), ly(0.62), lx(0.44), ly(0.08), s["subtitle"], 17, False, GRAY)
        else:
            rect(slide, 0, 0, SW, SH, PRIMARY)
            deco_circles(slide)
            rect(slide, 0, ly(0.86), SW, ly(0.14), DARK, alpha=60)
            rect(slide, lx(0.10), ly(0.315), lx(0.055), ly(0.011), ACCENT)
            text(slide, lx(0.10), ly(0.36), lx(0.8), ly(0.20), s["title"], 40, True, WHITE)
            if s.get("subtitle"):
                text(slide, lx(0.10), ly(0.56), lx(0.8), ly(0.08), s["subtitle"], 18, False,
                     C("F2DCDB"))
        return slide

    def r_toc(s):
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"] or "目录")
        items = s["points"][:5]
        v = layouts.get("toc", "numbered_list")
        if v == "cards":
            cols = 2
            cw, ch = 0.415, 0.155
            for i, it in enumerate(items):
                cx = 0.065 + (i % cols) * (cw + 0.035)
                cy = 0.28 + (i // cols) * (ch + 0.055)
                rect(slide, lx(cx), ly(cy), lx(cw), ly(ch), LIGHT, round_=True)
                rect(slide, lx(cx), ly(cy), lx(0.010), ly(ch), PRIMARY)
                text(slide, lx(cx + 0.03), ly(cy + 0.045), lx(0.07), ly(0.08), f"{i+1:02d}",
                     22, True, PRIMARY)
                text(slide, lx(cx + 0.105), ly(cy + 0.055), lx(cw - 0.13), ly(0.08), it,
                     17, True, DARK)
        else:
            y = 0.27
            for i, it in enumerate(items):
                chip = rect(slide, lx(0.07), ly(y), lx(0.042), lx(0.042), PRIMARY, round_=True)
                tf = chip.text_frame
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                r = para.add_run()
                r.text = f"{i+1:02d}"
                _font(r, 14, True, WHITE)
                text(slide, lx(0.135), ly(y + 0.006), lx(0.75), ly(0.06), it, 19, True, DARK)
                if i < len(items) - 1:
                    rect(slide, lx(0.07), ly(y + 0.088), lx(0.82), Emu(9525), LINE)
                y += 0.128
        return slide

    def r_section(s):
        slide = prs.slides.add_slide(blank)
        if layouts.get("section") == "center":
            rect(slide, 0, 0, SW, SH, LIGHT)
            circle(slide, SW * 0.9, SH * 0.15, SH * 0.25, PRIMARY, 10)
            circle(slide, SW * 0.08, SH * 0.9, SH * 0.2, PRIMARY, 8)
            rect(slide, lx(0.43), ly(0.40), lx(0.14), ly(0.009), ACCENT)
            text(slide, 0, ly(0.45), SW, ly(0.16), s["title"], 36, True, DARK,
                 align=PP_ALIGN.CENTER)
        else:
            rect(slide, 0, 0, lx(0.34), SH, PRIMARY)
            circle(slide, SW * 0.28, SH * 0.2, SH * 0.2, ACCENT, 16)
            circle(slide, SW * 0.06, SH * 0.85, SH * 0.24, ACCENT, 10)
            rect(slide, lx(0.34), 0, lx(0.008), SH, ACCENT)
            rect(slide, lx(0.06), ly(0.40), lx(0.05), ly(0.010), ACCENT)
            text(slide, lx(0.06), ly(0.44), lx(0.26), ly(0.18), s["title"], 30, True, WHITE)
        return slide

    def r_content(s):
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"])
        img = s.get("_image")
        has_img = bool(img and img.get("file_path"))
        if has_img:
            tw = 0.40
            y = 0.26 if layouts.get("content") == "bar_title" else 0.22
            for p in s["points"]:
                rect(slide, lx(0.06), ly(y + 0.014), lx(0.013), lx(0.013), ACCENT)
                text(slide, lx(0.085), ly(y), lx(tw), ly(0.13), p, 16, False, GRAY)
                y += 0.125 if len(s["points"]) <= 4 else 0.105
            ix, iy, iw, ih = lx(0.53), ly(0.235), lx(0.43), ly(0.60)
            rect(slide, ix + lx(0.008), iy + ly(0.015), iw, ih, LIGHT, round_=True)
            picture_fit(slide, img["file_path"], ix, iy, iw, ih)
            text(slide, ix, iy + ih + ly(0.012), iw, ly(0.045),
                 f"▲ {img.get('caption') or img['name']}", 11, False, GRAY, align=PP_ALIGN.CENTER)
        else:
            y0 = 0.26 if layouts.get("content") == "bar_title" else 0.22
            numbered_points(slide, s["points"], y=y0, size=17)
        return slide

    def r_data(s):
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"])
        cards = ((s.get("blocks") or {}).get("cards") or [])[:4]
        if cards:
            n = len(cards)
            cw = 0.88 / n - 0.02
            for i, c in enumerate(cards):
                cx = 0.06 + i * (cw + 0.02)
                rect(slide, lx(cx), ly(0.33), lx(cw), ly(0.30), LIGHT, round_=True)
                rect(slide, lx(cx), ly(0.33), lx(cw), ly(0.014), PRIMARY)
                text(slide, lx(cx), ly(0.395), lx(cw), ly(0.12), str(c.get("value", "")),
                     32, True, PRIMARY, align=PP_ALIGN.CENTER)
                rect(slide, lx(cx + cw / 2 - 0.02), ly(0.535), lx(0.04), ly(0.006), ACCENT)
                text(slide, lx(cx), ly(0.555), lx(cw), ly(0.06), str(c.get("label", "")),
                     14, False, GRAY, align=PP_ALIGN.CENTER)
        for j, p in enumerate(s["points"][:3]):
            rect(slide, lx(0.065), ly(0.735 + j * 0.068), lx(0.010), lx(0.010), ACCENT)
            text(slide, lx(0.088), ly(0.72 + j * 0.068), lx(0.84), ly(0.06), p, 14, False, GRAY)
        return slide

    def r_chart(s):
        """图表页：原生 pptx 图表（导出后可继续编辑数据）。"""
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"])
        chart_spec = (s.get("blocks") or {}).get("chart") or {}
        cats = chart_spec.get("categories") or []
        series = chart_spec.get("series") or []
        if cats and series:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            ct = (chart_spec.get("chart_type") or "bar").lower()
            chart_type = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                          "line": XL_CHART_TYPE.LINE_MARKERS,
                          "pie": XL_CHART_TYPE.PIE}.get(ct, XL_CHART_TYPE.COLUMN_CLUSTERED)
            data = CategoryChartData()
            data.categories = [str(c) for c in cats]
            n_series = 0
            for se in series:
                try:
                    data.add_series(str(se.get("name", "")), [float(v) for v in se.get("values", [])])
                    n_series += 1
                except (TypeError, ValueError):
                    continue
            if n_series:
                rect(slide, lx(0.055), ly(0.215), lx(0.89), ly(0.68), LIGHT, round_=True)
                gf = slide.shapes.add_chart(chart_type, lx(0.09), ly(0.25), lx(0.82), ly(0.60), data)
                gf.chart.has_legend = n_series > 1
                try:
                    palette = [PRIMARY, ACCENT, DARK, C("8FAADC")]
                    for si, serie in enumerate(gf.chart.series):
                        serie.format.fill.solid()
                        serie.format.fill.fore_color.rgb = palette[si % len(palette)]
                        serie.format.line.color.rgb = palette[si % len(palette)]
                except Exception:
                    pass
        for j, p in enumerate(s["points"][:2]):
            text(slide, lx(0.07), ly(0.90 + j * 0.045), lx(0.86), ly(0.045), p, 12, False, GRAY)
        return slide

    def r_case(s):
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"])
        rect(slide, lx(0.055), ly(0.235), lx(0.89), ly(0.64), LIGHT, round_=True)
        rect(slide, lx(0.055), ly(0.235), lx(0.012), ly(0.64), ACCENT)
        y = 0.30
        for p in s["points"]:
            text(slide, lx(0.10), ly(y), lx(0.80), ly(0.14), p, 17, False, DARK)
            y += 0.145
        return slide

    def r_timeline(s):
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"])
        items = ((s.get("blocks") or {}).get("timeline") or [])[:5]
        n = max(len(items), 1)
        yline = 0.52
        rect(slide, lx(0.08), ly(yline), lx(0.84), ly(0.009), ACCENT, round_=True)
        for i, it in enumerate(items):
            cx = 0.10 + i * (0.80 / max(n - 1, 1))
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx(cx - 0.009), ly(yline - 0.014),
                                         lx(0.018), lx(0.018))
            dot.fill.solid()
            dot.fill.fore_color.rgb = PRIMARY
            dot.line.color.rgb = WHITE
            dot.line.width = Pt(2)
            up = (i % 2 == 0)
            text(slide, lx(cx - 0.07), ly(yline - 0.105 if up else yline + 0.035), lx(0.14),
                 ly(0.05), str(it.get("time", "")), 14, True, PRIMARY, align=PP_ALIGN.CENTER)
            text(slide, lx(cx - 0.07), ly(yline - 0.06 if up else yline + 0.08), lx(0.14),
                 ly(0.10), str(it.get("text", "")), 12, False, GRAY, align=PP_ALIGN.CENTER)
        for j, p in enumerate(s["points"][:2]):
            text(slide, lx(0.07), ly(0.80 + j * 0.06), lx(0.86), ly(0.05), p, 14, False, GRAY)
        return slide

    def r_process(s):
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"])
        steps = ((s.get("blocks") or {}).get("process") or s["points"])[:5]
        n = max(len(steps), 1)
        cw = 0.88 / n
        for i, st in enumerate(steps):
            cx = 0.06 + i * cw
            shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, lx(cx), ly(0.36), lx(cw - 0.008),
                                         ly(0.13))
            shp.fill.solid()
            shp.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else DARK
            shp.line.fill.background()
            shp.shadow.inherit = False
            tf = shp.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(st)
            _font(r, 13, True, WHITE)
            text(slide, lx(cx), ly(0.53), lx(cw - 0.008), ly(0.07), f"第{i+1}步", 12, False,
                 GRAY, align=PP_ALIGN.CENTER)
        for j, p in enumerate(s["points"][:3]):
            text(slide, lx(0.07), ly(0.70 + j * 0.06), lx(0.86), ly(0.05), p, 14, False, GRAY)
        return slide

    def r_summary(s):
        slide = prs.slides.add_slide(blank)
        page_header(slide, s["title"])
        pts = s["points"][:4]
        cw = 0.425
        for i, p in enumerate(pts):
            cx = 0.06 + (i % 2) * (cw + 0.04)
            cy = 0.26 + (i // 2) * 0.25
            rect(slide, lx(cx), ly(cy), lx(cw), ly(0.20), LIGHT, round_=True)
            rect(slide, lx(cx), ly(cy), lx(0.012), ly(0.20), ACCENT)
            text(slide, lx(cx + 0.032), ly(cy + 0.035), lx(cw - 0.06), ly(0.15), p, 15, False, DARK)
        return slide

    def r_closing(s):
        slide = prs.slides.add_slide(blank)
        rect(slide, 0, 0, SW, SH, PRIMARY)
        deco_circles(slide)
        if layouts.get("closing") == "brand_band":
            rect(slide, 0, ly(0.86), SW, ly(0.14), DARK, alpha=60)
        rect(slide, lx(0.44), ly(0.50), lx(0.12), ly(0.010), ACCENT)
        text(slide, 0, ly(0.40), SW, ly(0.14), s["title"] or "谢谢聆听", 42, True, WHITE,
             align=PP_ALIGN.CENTER)
        return slide

    renderers = {"cover": r_cover, "toc": r_toc, "section": r_section, "content": r_content,
                 "data": r_data, "chart": r_chart, "case": r_case, "timeline": r_timeline,
                 "process": r_process, "summary": r_summary, "closing": r_closing}

    slides_list = outline.get("slides", [])
    total = len(slides_list)
    for idx, s in enumerate(slides_list, 1):
        slide = renderers.get(s.get("type"), r_content)(s)
        if s.get("type") in ("content", "toc", "data", "chart", "case", "timeline",
                             "process", "summary"):
            footer(slide, idx, total)
        notes(slide, s.get("note"))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ================= AI 页面级操作（编辑器右侧 AI 助手用） =================

_SLIDE_SYSTEM = "你是 PPT 内容编辑模块，只输出 JSON，不要用 markdown 包裹。"


def _slide_json(slide: Dict) -> str:
    return json.dumps({k: v for k, v in slide.items() if not k.startswith("_")},
                      ensure_ascii=False)


def ai_slide_action(assistant, action: str, slide: Dict, instruction: str = "") -> Dict:
    """AI 重写 / 扩写 / 精简 / 自定义指令修改单页。返回更新后的 slide。"""
    action_desc = {
        "rewrite": "整体重写这一页：换更清晰有力的表达，保持事实不变",
        "expand": "扩写这一页：把要点展开为更充实的表述，每条不超过40字",
        "condense": "精简这一页：压缩要点数量与字数，保留核心信息",
        "custom": instruction or "优化这一页",
    }.get(action, instruction or "优化这一页")

    prompt = f"""请处理一页 PPT 内容。

【操作要求】{action_desc}
【页面类型】{slide.get('type')}
【当前页面 JSON】
{_slide_json(slide)}

【输出要求】
1. 只输出修改后的该页 JSON（保持相同字段结构，保留 id/type/layout/blocks 结构，只改文字内容）；
2. 要点仍是 15~40 字的完整句子，有信息量；
3. 不要虚构事实和数据。"""
    raw = assistant.complete(
        [{"role": "system", "content": _SLIDE_SYSTEM},
         {"role": "user", "content": prompt}],
        temperature=0.5, max_tokens=1500)
    m = re.search(r"\{.*\}", raw or "", re.S)
    if not m:
        raise RuntimeError("AI 未返回有效结果，请重试")
    updated = json.loads(m.group(0))
    merged = {**slide, **updated, "id": slide.get("id") or uuid.uuid4().hex[:8]}
    merged.pop("_image", None)
    return normalize_outline({"slides": [merged]})["slides"][0]


def ai_generate_visual(assistant, kind: str, slide: Dict, instruction: str = "") -> Dict:
    """AI 生成图表 / 时间轴 / 流程图 / 数据卡片，写入 slide.blocks。

    kind: chart / timeline / process / data
    """
    spec_desc = {
        "chart": '{"chart":{"chart_type":"bar或pie或line","categories":["类别"...],"series":[{"name":"系列","values":[数值...]}]}}',
        "timeline": '{"timeline":[{"time":"时间点","text":"事项"},...]}（3~5 个节点）',
        "process": '{"process":["步骤1","步骤2",...]}（3~5 步）',
        "data": '{"cards":[{"label":"指标名","value":"数值"},...]}（3~4 张卡）',
    }
    prompt = f"""请根据页面内容生成"{kind}"可视化数据。

【页面内容】
{_slide_json(slide)}
【补充要求】{instruction or '无'}

【输出要求】
1. 只输出 JSON，格式：{{"blocks": {spec_desc[kind]}}}；
2. 数据必须来自页面内容或补充要求，不要虚构数字；
3. 页面内容不足以生成时，返回 {{"error": "原因"}}。"""
    raw = assistant.complete(
        [{"role": "system", "content": _SLIDE_SYSTEM},
         {"role": "user", "content": prompt}],
        temperature=0.4, max_tokens=1200)
    m = re.search(r"\{.*\}", raw or "", re.S)
    if not m:
        raise RuntimeError("AI 未返回有效结果，请重试")
    result = json.loads(m.group(0))
    if result.get("error"):
        raise RuntimeError(result["error"])
    updated = dict(slide)
    updated["type"] = kind
    updated["blocks"] = result.get("blocks") or {}
    return updated


def ai_structure_action(assistant, action: str, slides: List[Dict], index: int,
                        instruction: str = "") -> List[Dict]:
    """AI 增加 / 拆分 / 合并页面，返回新的 slides 数组。

    action: add（在 index 后新增一页）/ split（把 index 页拆成两页）/ merge（合并 index 与下一页）
    """
    ctx = json.dumps([{k: v for k, v in s.items() if not k.startswith("_")} for s in slides],
                     ensure_ascii=False)[:3000]
    target = slides[index] if 0 <= index < len(slides) else None

    if action == "add":
        prompt = f"""这份 PPT 的大纲如下：
{ctx}
用户要求：在第 {index + 1} 页之后新增一页。{instruction or '内容与上下文衔接，补充有价值的信息。'}
只输出新增页面的 JSON（type/title/points/blocks 结构，内容要实，不要虚构数据）。"""
        raw = assistant.complete([{"role": "system", "content": _SLIDE_SYSTEM},
                                  {"role": "user", "content": prompt}],
                                 temperature=0.5, max_tokens=1200)
        m = re.search(r"\{.*\}", raw or "", re.S)
        if not m:
            raise RuntimeError("AI 未返回有效结果")
        new_slide = json.loads(m.group(0))
        new_slide["id"] = uuid.uuid4().hex[:8]
        new_slide.setdefault("type", "content")
        new_slide.setdefault("points", [])
        new_slide.setdefault("blocks", {})
        slides = slides[:index + 1] + [new_slide] + slides[index + 1:]

    elif action == "split":
        if not target:
            raise RuntimeError("目标页面不存在")
        pts = target.get("points") or []
        if len(pts) < 2:
            raise RuntimeError("该页内容太少，无法拆分")
        half = (len(pts) + 1) // 2
        a = dict(target, points=pts[:half])
        b = dict(target, id=uuid.uuid4().hex[:8],
                 title=(target.get("title") or "") + "（续）", points=pts[half:])
        slides = slides[:index] + [a, b] + slides[index + 1:]

    elif action == "merge":
        if not target or index + 1 >= len(slides):
            raise RuntimeError("没有可合并的下一页")
        nxt = slides[index + 1]
        merged = dict(target)
        merged["points"] = (target.get("points") or []) + (nxt.get("points") or [])
        merged["blocks"] = target.get("blocks") or nxt.get("blocks") or {}
        slides = slides[:index] + [merged] + slides[index + 2:]
    else:
        raise RuntimeError("不支持的操作")

    return normalize_outline({"title": "", "slides": slides})["slides"]


# ================= 上传模板：从 pptx 提取视觉风格 =================

def _lighten(hex_color: str, factor: float = 0.92) -> str:
    """把颜色向白色混合，生成浅底色。"""
    r, g, b = (int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
    f = lambda v: int(v + (255 - v) * factor)
    return f"{f(r):02X}{f(g):02X}{f(b):02X}"


def extract_style_from_pptx(file_bytes: bytes) -> Dict:
    """解析上传的 pptx，提取完整主题色板（accent1-6）、母版背景色和字体。

    颜色映射：primary=accent1，accent=accent2（无则 accent3/1），
    dark=dk2，light=lt2 或 primary 的浅色衍生。
    """
    import zipfile
    colors = {}
    font = "微软雅黑"
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            theme_files = sorted(n for n in z.namelist() if n.startswith("ppt/theme/theme"))
            if theme_files:
                xml = z.read(theme_files[0]).decode("utf-8", "ignore")

                def _clr(tag, xml_text):
                    m = re.search(rf'<a:{tag}>.*?val="([0-9A-Fa-f]{{6}})"', xml_text, re.S)
                    if not m:
                        m = re.search(rf'<a:{tag}>.*?lastClr="([0-9A-Fa-f]{{6}})"', xml_text, re.S)
                    return m.group(1).upper() if m else None

                accents = [_clr(f"accent{i}", xml) for i in range(1, 7)]
                accents = [a for a in accents if a]
                if accents:
                    colors["primary"] = accents[0]
                    colors["accent"] = accents[1] if len(accents) > 1 else accents[0]
                dk2 = _clr("dk2", xml)
                if dk2:
                    colors["dark"] = dk2
                m = re.search(r'<a:ea typeface="([^"]+)"', xml)
                if m and m.group(1):
                    font = m.group(1)
                else:
                    m = re.search(r'<a:latin typeface="([^"]+)"', xml)
                    if m and m.group(1) and m.group(1) != "+mj-lt":
                        font = m.group(1)
            masters = sorted(n for n in z.namelist()
                             if re.match(r"ppt/slideMasters/slideMaster\d+\.xml$", n))
            for mf in masters:
                mxml = z.read(mf).decode("utf-8", "ignore")
                m = re.search(r'<p:bg>.*?val="([0-9A-Fa-f]{6})"', mxml, re.S)
                if m:
                    bg = m.group(1).upper()
                    if bg != "FFFFFF":
                        colors["light"] = bg
                    break
    except Exception as e:
        logger.warning("模板解析失败，使用默认风格: %s", e)
    if "primary" not in colors:
        colors = DEFAULT_COLORS.copy()
    else:
        base = DEFAULT_COLORS.copy()
        base.update(colors)
        if "light" not in base or base["light"] == DEFAULT_COLORS.get("light"):
            base["light"] = _lighten(base["primary"], 0.93)
        colors = base
    return {"colors": colors, "font": font, "layouts": DEFAULT_LAYOUTS.copy()}


# ================= 模板学习：分析上传模板的版式体系 =================

LAYOUT_TYPE_LABELS = {
    "cover": "封面", "toc": "目录", "section": "章节页", "content": "正文页",
    "content_image": "图文页", "two_col": "双栏页", "data": "数据页",
    "chart": "图表页", "case": "案例页", "timeline": "时间轴页",
    "process": "流程页", "summary": "总结页", "closing": "结束页",
}

# 大纲页型 → 版式库匹配的优先顺序（模板没有该版式时降级，不强行套用默认版式）
LEARNED_FALLBACK = {
    "cover": ["cover", "section", "content"],
    "toc": ["toc", "content", "two_col"],
    "section": ["section", "content", "cover"],
    "content": ["content", "content_image", "two_col", "section"],
    "data": ["data", "chart", "content", "two_col"],
    "chart": ["chart", "data", "content"],
    "case": ["case", "content_image", "content"],
    "timeline": ["timeline", "process", "content"],
    "process": ["process", "timeline", "content"],
    "summary": ["summary", "content", "section"],
    "closing": ["closing", "cover", "section", "content"],
}


def _shape_kind(shape) -> str:
    """形状分类：text/picture/table/chart/shape。"""
    try:
        if shape.shape_type == 13:  # PICTURE
            return "picture"
    except Exception:
        pass
    try:
        if shape.has_table:
            return "table"
    except Exception:
        pass
    try:
        if shape.has_chart:
            return "chart"
    except Exception:
        pass
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return "text"
    except Exception:
        pass
    return "shape"


def _region_of(shape, sw, sh) -> Dict:
    """形状位置归一化为页面比例。"""
    return {"x": round(shape.left / sw, 4), "y": round(shape.top / sh, 4),
            "w": round(shape.width / sw, 4), "h": round(shape.height / sh, 4)}


def _text_meta(shape) -> Dict:
    """文本框的排版特征：最大字号/是否加粗/颜色/对齐/示例文本。"""
    size, bold, color, align, sample = 0, False, None, None, ""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size.pt > size:
                    size = run.font.size.pt
                if run.font.bold:
                    bold = True
                if color is None:
                    try:
                        if run.font.color and run.font.color.type is not None:
                            color = str(run.font.color.rgb)
                    except Exception:
                        pass
            if align is None and para.alignment is not None:
                align = int(para.alignment)
        sample = shape.text_frame.text.strip()[:60]
    except Exception:
        pass
    return {"font_size": size or 18, "bold": bold, "color": color,
            "align": align, "sample": sample}


def _element_of(shape, sw, sh) -> Dict:
    """提取单个形状的结构信息；失败返回 None。"""
    try:
        kind = _shape_kind(shape)
        r = _region_of(shape, sw, sh)
        el = {"kind": kind, **r}
        if kind == "text":
            el.update(_text_meta(shape))
        elif kind == "picture":
            try:
                import hashlib
                el["hash"] = hashlib.md5(shape.image.blob).hexdigest()[:12]
            except Exception:
                el["hash"] = ""
        elif kind == "shape":
            try:
                el["shape_type"] = str(shape.auto_shape_type)
            except Exception:
                el["shape_type"] = ""
            try:
                if shape.fill.type is not None:
                    el["fill"] = str(shape.fill.fore_color.rgb)
            except Exception:
                pass
        return el
    except Exception:
        return None


def _slide_elements(slide, sw, sh) -> list:
    """提取一页的全部元素结构。"""
    els = []
    for shape in slide.shapes:
        el = _element_of(shape, sw, sh)
        if el:
            els.append(el)
    return els


def _classify_slide(els, idx, total) -> str:
    """规则分类页面版式。不依赖 LLM。"""
    texts = [e for e in els if e["kind"] == "text"]
    pics = [e for e in els if e["kind"] == "picture"]
    charts = [e for e in els if e["kind"] == "chart"]
    tables = [e for e in els if e["kind"] == "table"]
    shapes = [e for e in els if e["kind"] == "shape"]
    joined = " ".join(t.get("sample", "") for t in texts)
    max_size = max((t.get("font_size", 0) for t in texts), default=0)

    if idx == 0 and (max_size >= 28 or len(texts) <= 4):
        return "cover"
    if idx == total - 1 and re.search(r"谢|Thank|结束|致谢|汇报完毕", joined):
        return "closing"
    if re.search(r"目录|CONTENTS|Contents", joined) and len(texts) >= 3:
        return "toc"
    if charts:
        return "chart"
    if tables:
        return "data"
    arrows = [s for s in shapes if s.get("shape_type") and
              ("CHEVRON" in s["shape_type"] or "ARROW" in s["shape_type"])]
    if len(arrows) >= 3:
        return "process"
    if texts and len(texts) <= 3 and max_size >= 28 and len(joined) < 60 \
            and all(t.get("font_size", 0) >= max_size * 0.7 for t in texts):
        return "section"
    big_nums = [t for t in texts if t.get("font_size", 0) >= 24 and
                re.search(r"\d", t.get("sample", ""))]
    if len(big_nums) >= 2:
        return "data"
    if pics and texts:
        left = [t for t in texts if t["x"] < 0.45]
        right = [t for t in texts if t["x"] >= 0.45]
        if left and right and not pics:
            return "two_col"
        return "content_image"
    if texts:
        left = [t for t in texts if t["x"] < 0.45 and t["w"] < 0.5]
        right = [t for t in texts if t["x"] >= 0.45]
        if left and right:
            return "two_col"
    return "content"


def _layout_signature(els) -> tuple:
    """结构签名：用于去重。相同结构的页面只保留一个版式。"""
    sig = []
    for e in sorted(els, key=lambda e: (e["kind"], e["y"], e["x"])):
        sig.append((e["kind"], round(e["x"] * 6), round(e["y"] * 6),
                    round(e["w"] * 6), round(e["h"] * 6)))
    return tuple(sig)


def _mark_fixed(slides_els: list):
    """识别固定元素（logo/装饰/页脚）：在 >=60% 页面同一位置出现的形状/图片/相同文本。
    直接修改每个元素的 role 字段：fixed / content。"""
    from collections import Counter
    n = len(slides_els)
    counter = Counter()
    for els in slides_els:
        keys = set()
        for e in els:
            if e["kind"] == "picture":
                k = ("pic", e.get("hash", ""), round(e["x"] * 10), round(e["y"] * 10))
            elif e["kind"] == "shape":
                k = ("shape", round(e["x"] * 10), round(e["y"] * 10),
                     round(e["w"] * 10), round(e["h"] * 10), e.get("fill", ""))
            elif e["kind"] == "text" and len(e.get("sample", "")) <= 15:
                k = ("txt", e.get("sample", ""), round(e["x"] * 10), round(e["y"] * 10))
            else:
                continue
            keys.add(k)
        counter.update(keys)
    threshold = max(2, int(n * 0.6))
    for els in slides_els:
        for e in els:
            if e["kind"] == "picture":
                k = ("pic", e.get("hash", ""), round(e["x"] * 10), round(e["y"] * 10))
            elif e["kind"] == "shape":
                k = ("shape", round(e["x"] * 10), round(e["y"] * 10),
                     round(e["w"] * 10), round(e["h"] * 10), e.get("fill", ""))
            elif e["kind"] == "text" and len(e.get("sample", "")) <= 15:
                k = ("txt", e.get("sample", ""), round(e["x"] * 10), round(e["y"] * 10))
            else:
                k = None
            e["role"] = "fixed" if k and counter[k] >= threshold else "content"


def analyze_template(file_bytes: bytes) -> Dict:
    """模板学习主入口：解析上传的 pptx，自动建立版式库。

    返回 {slide_size, colors, font, layouts: [...]}；解析失败时 layouts 为空，
    调用方回退 DEFAULT_LAYOUTS。
    """
    from pptx import Presentation
    style = extract_style_from_pptx(file_bytes)
    result = {"colors": style["colors"], "font": style["font"],
              "slide_size": None, "layouts": []}
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.warning("模板学习失败：%s", e)
        return result
    sw, sh = prs.slide_width, prs.slide_height
    result["slide_size"] = {"w": int(sw), "h": int(sh)}

    slides = list(prs.slides)
    if not slides:
        return result
    slides_els = [_slide_elements(s, sw, sh) for s in slides]
    _mark_fixed(slides_els)

    background = None
    try:
        mbg = prs.slide_masters[0].background.fill
        if mbg.type is not None:
            background = str(mbg.fore_color.rgb)
    except Exception:
        pass

    seen, layouts, type_count = set(), [], {}
    total = len(slides)
    for idx, (slide, els) in enumerate(zip(slides, slides_els)):
        sig = _layout_signature([e for e in els if e.get("role") != "fixed"])
        if sig in seen:
            continue
        seen.add(sig)
        dtype = _classify_slide(els, idx, total)
        type_count[dtype] = type_count.get(dtype, 0) + 1
        base_name = LAYOUT_TYPE_LABELS.get(dtype, "内容页")
        name = base_name if type_count[dtype] == 1 else f"{base_name}{type_count[dtype]}"

        texts = [e for e in els if e["kind"] == "text" and e.get("role") != "fixed"]
        texts.sort(key=lambda t: (-t.get("font_size", 0), t["y"], t["x"]))
        text_regions = []
        for ti, t in enumerate(texts):
            text_regions.append({
                "role": "title" if ti == 0 else ("subtitle" if dtype == "cover" and ti == 1 else "body"),
                **{k: t[k] for k in ("x", "y", "w", "h")},
                "font_size": t.get("font_size", 18), "bold": t.get("bold", False),
                "color": t.get("color"), "align": t.get("align")})
        layouts.append({
            "id": f"L{idx}",
            "name": name,
            "detected_type": dtype,
            "source_slide_index": idx,
            "slide_size": {"w": int(sw), "h": int(sh)},
            "background": background,
            "element_schema": els,
            "placeholders": [{"idx": ph.placeholder_format.idx,
                              "type": str(ph.placeholder_format.type)}
                             for ph in slide.placeholders] if slide.placeholders else [],
            "text_regions": text_regions,
            "image_regions": [{k: e[k] for k in ("x", "y", "w", "h")}
                              for e in els if e["kind"] == "picture" and e.get("role") != "fixed"],
            "chart_regions": [{k: e[k] for k in ("x", "y", "w", "h")}
                              for e in els if e["kind"] in ("chart", "table") and e.get("role") != "fixed"],
            "shape_regions": [{k: e[k] for k in ("x", "y", "w", "h")}
                              for e in els if e["kind"] == "shape" and e.get("role") == "fixed"],
        })
    result["layouts"] = layouts
    return result


# ================= 模板学习渲染：按版式库克隆页面并填充内容 =================

def _pick_layout(layouts: list, slide_type: str, used: Dict) -> Dict:
    """按页型从版式库选最合适的版式；同类型多版式时轮流使用。无匹配则降级链查找。"""
    if not layouts:
        return None
    chain = LEARNED_FALLBACK.get(slide_type, ["content"]) + ["content"]
    for want in chain:
        cands = [l for l in layouts if l.get("detected_type") == want]
        if cands:
            i = used.get(want, 0) % len(cands)
            used[want] = used.get(want, 0) + 1
            return cands[i]
    return layouts[0]


def _clone_slide(prs, src_slide):
    """在同一演示文稿内克隆一页：深拷贝所有形状，并迁移图片关系。"""
    import copy
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    rid_map = {}
    for shp in src_slide.shapes:
        new_el = copy.deepcopy(shp._element)
        new_slide.shapes._spTree.append(new_el)
    from pptx.oxml.ns import qn as _qn
    for blip in new_slide.shapes._spTree.iter(_qn("a:blip")):
        old_rid = blip.get(_qn("r:embed"))
        if not old_rid:
            continue
        if old_rid not in rid_map:
            try:
                image_part = src_slide.part.rels[old_rid].target_part
                new_rid = new_slide.part.relate_to(image_part, RT.IMAGE)
                rid_map[old_rid] = new_rid
            except Exception:
                continue
        blip.set(_qn("r:embed"), rid_map[old_rid])
    return new_slide


def _fill_text_region(shape, lines, keep_para_format=True):
    """向文本框填充内容：用首个段落的格式作为模板逐条写入，清除多余示例段落。"""
    import copy
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    if not paras:
        return
    proto = paras[0]._p
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    if not lines:
        first.text = ""
        return
    first.text = str(lines[0])
    for line in lines[1:]:
        new_p = copy.deepcopy(proto)
        for r in new_p.findall(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}r"):
            new_p.remove(r)
        tf._txBody.append(new_p)
        tf.paragraphs[-1].text = str(line)


def _clear_text(shape):
    try:
        tf = shape.text_frame
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        if tf.paragraphs:
            tf.paragraphs[0].text = ""
    except Exception:
        pass


def _delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def _add_native_chart(slide, chart_spec, region, colors, font, sw, sh):
    """在模板图表区域绘制原生可编辑图表（替换示例图表）。"""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    C = lambda h: RGBColor.from_string(h)
    x, y = Emu(int(region["x"] * sw)), Emu(int(region["y"] * sh))
    w, h = Emu(int(region["w"] * sw)), Emu(int(region["h"] * sh))
    cd = CategoryChartData()
    cd.categories = chart_spec.get("categories", [])
    for serie in chart_spec.get("series", []):
        cd.add_series(serie.get("name", ""), serie.get("values", []))
    ct = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
          "line": XL_CHART_TYPE.LINE_MARKERS,
          "pie": XL_CHART_TYPE.PIE}.get(chart_spec.get("chart_type", "bar"),
                                        XL_CHART_TYPE.COLUMN_CLUSTERED)
    gf = slide.shapes.add_chart(ct, x, y, w, h, cd)
    chart = gf.chart
    chart.has_title = False
    try:
        chart.font.name = font
        chart.font.size = Pt(11)
    except Exception:
        pass
    if chart_spec.get("chart_type") != "pie" and len(chart_spec.get("series", [])) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    palette = [C(colors["primary"]), C(colors["accent"]), C(colors["dark"]), C("8FAADC")]
    for si, serie in enumerate(chart.series):
        serie.format.fill.solid()
        serie.format.fill.fore_color.rgb = palette[si % len(palette)]
    return gf


def _fill_learned_slide(slide, layout, sdata, colors, font, sw, sh, prs):
    """把大纲一页的内容填进克隆出来的模板页：
    - 标题/副标题/正文 → 对应文本区域（示例文字全部替换）
    - 示例图片/图表/表格（非固定元素）→ 删除；有匹配素材图时在图片区插入
    - 固定元素（logo/装饰/页脚）→ 原样保留
    - 有图表块且有图表区域 → 绘制原生图表
    """
    schema = layout.get("element_schema", [])

    def role_of(e):
        for se in schema:
            if se["kind"] == e["kind"] and abs(se["x"] - e["x"]) < 0.01 \
                    and abs(se["y"] - e["y"]) < 0.01:
                return se.get("role", "content")
        return "content"

    texts, pics, chartframes = [], [], []
    for shp in list(slide.shapes):
        e = _element_of(shp, sw, sh)
        if not e:
            continue
        role = role_of(e)
        kind = e["kind"]
        if role == "fixed":
            continue
        if kind == "text":
            texts.append((shp, e))
        elif kind == "picture":
            pics.append((shp, e))
        elif kind in ("chart", "table"):
            chartframes.append((shp, e))

    texts.sort(key=lambda te: (-te[1].get("font_size", 0), te[1]["y"], te[1]["x"]))
    title = sdata.get("title", "")
    subtitle = sdata.get("subtitle", "")
    points = [str(p) for p in (sdata.get("points") or [])]
    blocks = sdata.get("blocks") or {}

    if blocks.get("cards") and not layout.get("chart_regions"):
        for c in blocks["cards"]:
            points.append(f"{c.get('label', '')}：{c.get('value', '')}")
    if blocks.get("timeline"):
        for t in blocks["timeline"]:
            points.append(f"{t.get('time', '')}　{t.get('event', '')}")
    if blocks.get("process"):
        for i, p in enumerate(blocks["process"], 1):
            points.append(f"第{i}步：{p if isinstance(p, str) else p.get('name', '')}")

    is_cover = layout.get("detected_type") == "cover"
    body_filled = False
    for ti, (shp, e) in enumerate(texts):
        if ti == 0:
            _fill_text_region(shp, [title] if title else [])
        elif is_cover and ti == 1:
            _fill_text_region(shp, [subtitle] if subtitle else [])
        elif not body_filled and points:
            _fill_text_region(shp, points)
            body_filled = True
        else:
            _clear_text(shp)

    img_path = (sdata.get("_image") or {}).get("file_path")
    for pi, (shp, e) in enumerate(pics):
        if pi == 0 and img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, shp.left, shp.top, shp.width, shp.height)
            except Exception:
                pass
        _delete_shape(shp)

    chart_spec = blocks.get("chart")
    chart_drawn = False
    for shp, e in chartframes:
        if chart_spec and not chart_drawn:
            region = {k: e[k] for k in ("x", "y", "w", "h")}
            _delete_shape(shp)
            try:
                _add_native_chart(slide, chart_spec, region, colors, font, sw, sh)
                chart_drawn = True
            except Exception as ex:
                logger.warning("图表绘制失败: %s", ex)
        else:
            _delete_shape(shp)
    if chart_spec and not chart_drawn:
        names = "、".join(sr.get("name", "") for sr in chart_spec.get("series", []))
        if texts and len(texts) > 1:
            shp, e = texts[1]
            _fill_text_region(shp, points + [f"（图表数据：{names}）"] if points else [f"（图表数据：{names}）"])


def build_pptx_learned(outline: Dict, template: Dict, base_pptx: bytes,
                       layout_library: list) -> bytes:
    """模板学习渲染：克隆模板页 → 清除示例内容 → 填充新内容。
    模板没有的页型自动降级到已有版式，不套用 DEFAULT_LAYOUTS。"""
    from pptx import Presentation
    colors = template["colors"]
    font = template.get("font") or "微软雅黑"
    prs = Presentation(io.BytesIO(base_pptx))
    src_slides = list(prs.slides)
    used = {}
    plan = []
    for s in outline.get("slides", []):
        lay = _pick_layout(layout_library, s.get("type", "content"), used)
        if not lay:
            break
        idx = lay.get("source_slide_index", 0)
        if idx >= len(src_slides):
            idx = 0
        plan.append((s, lay, src_slides[idx]))
    new_slides = []
    for s, lay, src in plan:
        try:
            ns = _clone_slide(prs, src)
            new_slides.append((ns, lay, s))
        except Exception as e:
            logger.warning("克隆页面失败: %s", e)
    from pptx.oxml.ns import qn as _qn2
    sldIdLst = prs.slides._sldIdLst
    for sld in list(sldIdLst)[:len(src_slides)]:
        rid = sld.get(_qn2('r:id'))
        sldIdLst.remove(sld)
        try:
            prs.part.drop_rel(rid)
        except Exception:
            pass
    sw, sh = prs.slide_width, prs.slide_height
    for ns, lay, s in new_slides:
        try:
            _fill_learned_slide(ns, lay, s, colors, font, sw, sh, prs)
        except Exception as e:
            logger.warning("页面填充失败: %s", e)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_pptx(outline: Dict, template: Dict, base_pptx: bytes = None) -> bytes:
    """渲染总入口：有版式库+底版 → 模板学习渲染；否则 → 精修渲染（fallback）。"""
    library = template.get("layout_library") or []
    if base_pptx and library:
        try:
            return build_pptx_learned(outline, template, base_pptx, library)
        except Exception as e:
            logger.warning("模板学习渲染失败，回退默认渲染: %s", e)
    return build_pptx(outline, template, base_pptx=base_pptx)
